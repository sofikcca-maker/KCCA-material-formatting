"""
Doc-to-Deck: Extract content from a source document, structure it with Claude,
and populate a user-supplied PowerPoint template — preserving the template's
exact fonts, colors, and layout.

Run with: streamlit run app.py
"""

import io
import json
import re

import streamlit as st
from pptx import Presentation

# Optional/format-specific parsers are imported lazily inside the functions
# that need them, so the app still runs if one library is missing and the
# user never uploads that file type.


# --------------------------------------------------------------------------
# Document text extraction
# --------------------------------------------------------------------------

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract plain text from a PDF using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages_text = []
    for page in reader.pages:
        text = page.extract_text() or ""
        pages_text.append(text)
    return "\n".join(pages_text)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract plain text from a Word document using python-docx."""
    import docx  # python-docx

    document = docx.Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    # Also pull text out of tables, since reports often use them.
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text)
    return "\n".join(paragraphs)


def extract_text_from_plain(file_bytes: bytes) -> str:
    """Decode plain text / markdown files, tolerating odd encodings."""
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    # Last resort: replace undecodable bytes rather than crashing.
    return file_bytes.decode("utf-8", errors="replace")


def extract_document_text(uploaded_file) -> str:
    """Dispatch to the correct extractor based on file extension."""
    name = uploaded_file.name.lower()
    file_bytes = uploaded_file.getvalue()

    if not file_bytes:
        raise ValueError("The uploaded source document is empty.")

    if name.endswith(".pdf"):
        text = extract_text_from_pdf(file_bytes)
    elif name.endswith(".docx"):
        text = extract_text_from_docx(file_bytes)
    elif name.endswith(".txt") or name.endswith(".md"):
        text = extract_text_from_plain(file_bytes)
    else:
        raise ValueError(
            f"Unsupported file format: '{uploaded_file.name}'. "
            "Please upload a .pdf, .docx, .txt, or .md file."
        )

    text = text.strip()
    if not text:
        raise ValueError(
            "No readable text could be extracted from the document "
            "(it may be a scanned/image-only PDF or an empty file)."
        )
    return text


# --------------------------------------------------------------------------
# Claude API: structure raw text into a slide outline
# --------------------------------------------------------------------------

CLAUDE_MODEL = "claude-3-5-sonnet-latest"

SYSTEM_PROMPT = """You are a presentation content strategist. You convert raw \
source documents into concise, well-structured slide outlines.

Analyze the text provided by the user and structure it into presentation \
slides. Return EXCLUSIVELY a JSON array adhering to this exact structure, \
with no markdown formatting, no code fences, and no introductory or \
trailing text:

[
  {
    "title": "Slide Title",
    "bullet_points": ["Key point 1", "Key point 2", "Key point 3"]
  }
]

Rules:
- Produce between 4 and 12 slides, depending on how much substantive \
content is in the source text.
- Each title should be short (under 8 words).
- Each slide should have 3 to 5 bullet points, each a single concise \
sentence or fragment (under ~18 words).
- Do not invent facts that are not supported by the source text.
- Output raw JSON only. Nothing else.
"""


def build_user_prompt(document_text: str, goal: str) -> str:
    """Compose the user-turn prompt sent to Claude."""
    goal_line = f"Presentation goal: {goal}\n\n" if goal else ""
    return (
        f"{goal_line}"
        "Source document text:\n"
        "-----\n"
        f"{document_text}\n"
        "-----\n\n"
        "Structure this into a slide outline following the required JSON "
        "format."
    )


def clean_json_response(raw_text: str) -> str:
    """Strip common wrapping artifacts (code fences, stray prose) from a
    model response so it can be parsed as JSON."""
    text = raw_text.strip()

    # Remove ```json ... ``` or ``` ... ``` code fences if present.
    fence_match = re.search(r"```(?:json)?\s*(.*?)\s*```", text, re.DOTALL)
    if fence_match:
        text = fence_match.group(1).strip()

    # If there's still leading/trailing prose, isolate the outermost
    # JSON array by locating the first '[' and the last ']'.
    start = text.find("[")
    end = text.rfind("]")
    if start != -1 and end != -1 and end > start:
        text = text[start : end + 1]

    return text.strip()


def generate_slide_outline(api_key: str, document_text: str, goal: str) -> list:
    """Call the Claude API and return a validated list of slide dicts."""
    import anthropic

    if not api_key or not api_key.strip():
        raise ValueError("Please enter your Claude API key in the sidebar.")

    client = anthropic.Anthropic(api_key=api_key.strip())

    # Guard against extremely long documents blowing the context window.
    max_chars = 100_000
    trimmed_text = document_text[:max_chars]

    try:
        response = client.messages.create(
            model=CLAUDE_MODEL,
            max_tokens=4096,
            system=SYSTEM_PROMPT,
            messages=[
                {
                    "role": "user",
                    "content": build_user_prompt(trimmed_text, goal),
                }
            ],
        )
    except anthropic.AuthenticationError as exc:
        raise ValueError(
            "Authentication failed — please check that your Claude API key "
            "is correct."
        ) from exc
    except anthropic.APIError as exc:
        raise ValueError(f"The Claude API returned an error: {exc}") from exc

    raw_text = "".join(
        block.text for block in response.content if block.type == "text"
    )

    cleaned = clean_json_response(raw_text)

    try:
        slides = json.loads(cleaned)
    except json.JSONDecodeError as exc:
        raise ValueError(
            "Claude's response could not be parsed as JSON. "
            "Try again, or try a shorter/cleaner source document."
        ) from exc

    if not isinstance(slides, list) or not slides:
        raise ValueError("Claude did not return a usable slide outline.")

    # Normalize/validate each slide entry.
    normalized = []
    for item in slides:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title", "")).strip()
        bullets = item.get("bullet_points", [])
        if not isinstance(bullets, list):
            bullets = [str(bullets)]
        bullets = [str(b).strip() for b in bullets if str(b).strip()]
        if title:
            normalized.append({"title": title, "bullet_points": bullets})

    if not normalized:
        raise ValueError("Claude's response did not contain any valid slides.")

    return normalized


# --------------------------------------------------------------------------
# PowerPoint population — preserves the uploaded template's styling
# --------------------------------------------------------------------------

def pick_content_layout(prs: Presentation):
    """Pick the best 'Title and Content' style layout from the template.

    Preference order:
      1. A layout whose name suggests title + body content.
      2. Layout index 1 (the conventional 'Title and Content' slot).
      3. Any layout that has at least a title and one other placeholder.
      4. The first available layout, as a last resort.
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("The uploaded template contains no slide layouts.")

    name_keywords = ("title and content", "title, content", "content")
    for layout in layouts:
        if layout.name and layout.name.strip().lower() in name_keywords:
            return layout

    if len(layouts) > 1:
        return layouts[1]

    return layouts[0]


def get_body_placeholder(slide):
    """Find the placeholder on the slide meant for bullet/body content,
    i.e. any placeholder that isn't the title."""
    from pptx.util import Emu  # noqa: F401  (kept for clarity/future use)

    title_shape = slide.shapes.title
    for placeholder in slide.placeholders:
        if title_shape is not None and placeholder.shape_id == title_shape.shape_id:
            continue
        if placeholder.has_text_frame:
            return placeholder
    return None


def populate_bullets(placeholder, bullet_points: list):
    """Write bullet text into a placeholder's text frame without touching
    fonts, sizes, colors, or paragraph-level styling defined by the
    template/layout. We only ever set the .text of paragraphs, which
    reuses the run/paragraph formatting already defined for that
    placeholder in the template."""
    text_frame = placeholder.text_frame

    if not bullet_points:
        text_frame.text = ""
        return

    # Setting .text on the first paragraph preserves that paragraph's
    # existing run formatting (font, size, color) as defined by the
    # template's placeholder/layout.
    text_frame.text = bullet_points[0]

    # Additional bullets reuse the placeholder's default paragraph
    # formatting (inherited from the layout/master) — we deliberately do
    # NOT set explicit font/size/color here, to preserve template fidelity.
    for bullet in bullet_points[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0


def build_presentation(template_bytes: bytes, slides_data: list) -> bytes:
    """Open the uploaded template in memory and append one slide per
    entry in slides_data, using the template's own layout and
    placeholder formatting throughout."""
    prs = Presentation(io.BytesIO(template_bytes))
    layout = pick_content_layout(prs)

    for slide_data in slides_data:
        slide = prs.slides.add_slide(layout)

        # Title placeholder
        if slide.shapes.title is not None:
            slide.shapes.title.text = slide_data["title"]
        else:
            # Layout has no title placeholder — fall back to the first
            # available text placeholder for the title text.
            fallback = next(
                (p for p in slide.placeholders if p.has_text_frame), None
            )
            if fallback is not None:
                fallback.text_frame.text = slide_data["title"]

        # Body / bullet placeholder
        body_placeholder = get_body_placeholder(slide)
        if body_placeholder is not None:
            populate_bullets(body_placeholder, slide_data["bullet_points"])
        # If the layout genuinely has no second placeholder, we simply
        # skip bullets for that slide rather than injecting a new text
        # box, to avoid deviating from the template's design.

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# --------------------------------------------------------------------------
# Streamlit UI
# --------------------------------------------------------------------------

def main():
    st.set_page_config(page_title="Doc-to-Deck", page_icon="📊", layout="centered")
    st.title("📊 Doc-to-Deck")
    st.caption(
        "Turn any document into a PowerPoint that matches your template's "
        "exact styling."
    )

    # ---------------- Sidebar ----------------
    with st.sidebar:
        st.header("Configuration")
        api_key = st.text_input(
            "Claude API Key",
            type="password",
            placeholder="sk-ant-...",
            help="Your key is used only for this session and is never stored.",
        )
        template_file = st.file_uploader(
            "Upload PowerPoint Template (.pptx)",
            type=["pptx"],
            help="This template's layout, fonts, and colors will be preserved exactly.",
        )

    # ---------------- Main panel ----------------
    st.subheader("1. Upload your source document")
    source_file = st.file_uploader(
        "Upload Source Document",
        type=["pdf", "docx", "txt", "md"],
    )

    st.subheader("2. Choose a presentation goal (optional)")
    goal_options = [
        "General Summary",
        "Executive Summary",
        "Sales Pitch",
        "Technical Report",
        "Project Update",
        "Custom...",
    ]
    goal_choice = st.selectbox("Presentation goal", goal_options, index=0)
    goal = goal_choice
    if goal_choice == "Custom...":
        goal = st.text_input("Describe the goal", placeholder="e.g. Investor pitch")

    st.subheader("3. Generate")
    generate_clicked = st.button("Generate PowerPoint", type="primary")

    status_placeholder = st.empty()

    if generate_clicked:
        # ---- Input validation ----
        if not api_key or not api_key.strip():
            status_placeholder.error("Please enter your Claude API key in the sidebar.")
            return
        if template_file is None:
            status_placeholder.error("Please upload a PowerPoint template (.pptx) in the sidebar.")
            return
        if source_file is None:
            status_placeholder.error("Please upload a source document to convert.")
            return

        try:
            with st.spinner("Extracting text from the source document..."):
                document_text = extract_document_text(source_file)

            with st.spinner("Asking Claude to structure the content into slides..."):
                slides_data = generate_slide_outline(api_key, document_text, goal)

            with st.spinner("Populating your PowerPoint template..."):
                template_bytes = template_file.getvalue()
                if not template_bytes:
                    raise ValueError("The uploaded template file is empty.")
                pptx_bytes = build_presentation(template_bytes, slides_data)

        except ValueError as exc:
            status_placeholder.error(str(exc))
            return
        except Exception as exc:  # noqa: BLE001 - surface any unexpected error to the user
            status_placeholder.error(f"An unexpected error occurred: {exc}")
            return

        status_placeholder.success(
            f"Done! Generated {len(slides_data)} slides while preserving your template's styling."
        )

        st.download_button(
            label="⬇️ Download Generated PowerPoint",
            data=pptx_bytes,
            file_name="generated_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        with st.expander("Preview extracted slide outline"):
            for i, slide in enumerate(slides_data, start=1):
                st.markdown(f"**Slide {i}: {slide['title']}**")
                for bullet in slide["bullet_points"]:
                    st.markdown(f"- {bullet}")


if __name__ == "__main__":
    main()